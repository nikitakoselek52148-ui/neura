from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from docx import Document
from pptx import Presentation
import google.generativeai as genai
import uuid
import os
import uvicorn

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-1.5-flash")

class UserRequest(BaseModel):
    text: str
    file_type: str = None

@app.post("/chat")
async def chat(req: UserRequest):
    try:
        response = model.generate_content(req.text)
        return {"response": response.text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/generate")
async def generate(req: UserRequest):
    try:
        response = model.generate_content(req.text)
        content = response.text

        if req.file_type == "docx":
            doc = Document()
            doc.add_heading("Документ", level=1)
            for line in content.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
            filename = f"doc_{uuid.uuid4().hex[:8]}.docx"
            doc.save(filename)
        else:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Презентация"
            body = slide.shapes.placeholders[1].text_frame
            for line in content.split('\n')[:30]:
                if line.strip():
                    body.add_paragraph(line)
            filename = f"pres_{uuid.uuid4().hex[:8]}.pptx"
            prs.save(filename)
        
        return FileResponse(filename, filename=filename)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/")
def root():
    return {"status": "ok", "model": "Gemini 1.5 Flash"}

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
